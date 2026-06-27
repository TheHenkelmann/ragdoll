# SPDX-License-Identifier: AGPL-3.0-only

from __future__ import annotations

from pathlib import Path

import httpx
import pytesseract
import trafilatura
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pdf2image import convert_from_path
from pptx import Presentation
from pypdf import PdfReader


def extract_from_source(source: dict[str, object], staging_dir: Path) -> str:
    source_type = str(source["type"])
    uri = source.get("uri")

    if source_type == "text":
        text_path = staging_dir / f"{source['id']}.txt"
        if text_path.exists():
            return text_path.read_text(encoding="utf-8", errors="ignore")
        raise RuntimeError("text staging file missing")

    if source_type == "url":
        if not uri:
            raise RuntimeError("url source missing uri")
        return extract_url(str(uri))

    if source_type == "file":
        if not uri:
            raise RuntimeError("file source missing uri")
        return extract_file(Path(str(uri)), name=str(source.get("name", "")))

    raise RuntimeError(f"unsupported source type: {source_type}")


def extract_url(url: str) -> str:
    response = httpx.get(url, timeout=30.0, follow_redirects=True)
    response.raise_for_status()
    downloaded = response.text
    extracted = trafilatura.extract(
        downloaded,
        include_comments=False,
        include_tables=True,
        output_format="markdown",
    )
    if extracted and extracted.strip():
        return extracted
    return downloaded


def extract_file(path: Path, *, name: str = "") -> str:
    suffix = path.suffix.lower()
    if not suffix and name:
        suffix = Path(name).suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix in {".xlsx", ".xlsm"}:
        return extract_xlsx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    raise RuntimeError(f"unsupported file type: {suffix}")


def build_pdf_page_map(path: Path) -> tuple[str, list[dict[str, int]]]:
    reader = PdfReader(str(path))
    parts: list[str] = []
    page_map: list[dict[str, int]] = []
    pos = 0
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if not text.strip():
            continue
        start = pos
        end = pos + len(text)
        page_map.append({"page": page_num, "start": start, "end": end})
        parts.append(text)
        pos = end + 2
    if parts:
        return "\n\n".join(parts), page_map
    ocr_text = ocr_pdf(path)
    return ocr_text, []


def extract_pdf(path: Path) -> str:
    text, _page_map = build_pdf_page_map(path)
    return text


def ocr_pdf(path: Path) -> str:
    images = convert_from_path(str(path))
    texts: list[str] = []
    for image in images:
        text = pytesseract.image_to_string(image, lang="deu+eng")
        if text.strip():
            texts.append(text)
    return "\n\n".join(texts)


def extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_xlsx(path: Path) -> str:
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                parts.append(" ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)
