# SPDX-License-Identifier: AGPL-3.0-only

from __future__ import annotations

import hashlib
from pathlib import Path

import numpy as np
import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from ragdoll_worker.config import WorkerConfig, sanitize_model_name
from ragdoll_worker.db import WorkerDb


@pytest.fixture
def data_dir(tmp_path: Path) -> Path:
    root = tmp_path / "data"
    (root / "db").mkdir(parents=True)
    (root / "staging").mkdir(parents=True)
    (root / "models").mkdir(parents=True)
    return root


@pytest.fixture
def worker_config(data_dir: Path) -> WorkerConfig:
    return WorkerConfig(
        data_dir=data_dir,
        db_path=data_dir / "db" / "ragdoll.db",
        model_cache_dir=data_dir / "models",
        staging_dir=data_dir / "staging",
        embedding_model="BAAI/bge-m3",
        embedding_dim=1024,
        worker_poll_interval_ms=1000,
        job_lease_seconds=60,
        max_attempts=3,
        worker_id="test-worker",
    )


@pytest.fixture
def worker_db(worker_config: WorkerConfig) -> WorkerDb:
    migrations = Path(__file__).resolve().parents[2] / "migrations" / "0001_init.sql"
    sql = migrations.read_text(encoding="utf-8")
    db = WorkerDb(worker_config)
    db.conn.executescript(sql)
    commit = getattr(db.conn, "commit", None)
    if callable(commit):
        commit()
    yield db
    db.release()


class MockTokenizer:
    def encode(self, text: str):
        return type("Enc", (), {"ids": list(range(max(1, len(text.split()))))})()


@pytest.fixture
def mock_tokenizer() -> MockTokenizer:
    return MockTokenizer()


class MockEmbedder:
    """Deterministic normalized vectors derived from each text's SHA-256 hash."""

    def __init__(self, dim: int = 8) -> None:
        self.dim = dim

    def _vector_for(self, text: str) -> np.ndarray:
        digest = hashlib.sha256(text.encode("utf-8")).digest()
        seed = int.from_bytes(digest[:8], "big")
        rng = np.random.default_rng(seed)
        vec = rng.standard_normal(self.dim)
        norm = float(np.linalg.norm(vec))
        if norm == 0.0:
            vec[0] = 1.0
            norm = 1.0
        return vec / norm

    def embed(self, texts: list[str]):
        for text in texts:
            yield self._vector_for(text)


@pytest.fixture
def mock_embedder() -> MockEmbedder:
    return MockEmbedder()


def make_docx(path: Path, paragraphs: list[str]) -> Path:
    doc = DocxDocument()
    for paragraph in paragraphs:
        doc.add_paragraph(paragraph)
    doc.save(str(path))
    return path


def make_xlsx(path: Path, rows: list[list[object]]) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    for row in rows:
        sheet.append(row)
    workbook.save(str(path))
    return path


def make_pptx(path: Path, slide_texts: list[str]) -> Path:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    for text in slide_texts:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        textbox.text_frame.text = text
    presentation.save(str(path))
    return path
